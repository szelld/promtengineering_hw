"""
PDF to PPTX Presentation Generator
===========================
Reads a PDF document, sends it to Google Gemini to extract key points
as structured JSON slides, and renders each slide as a separate slide
in a PPTX using python-pptx.

Usage:
    1. Ensure your Gemini API key is set.
    2. Run: python main.py
    3. Open Okosotthon1.pptx to view the generated presentation.
"""

import json
import sys
import os
from dotenv import load_dotenv

# ---------------------------------------------------------------------------
# Third-party imports
# ---------------------------------------------------------------------------
from google import genai
from google.genai import types
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from pypdf import PdfReader

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

# Load environment variables from .env file
load_dotenv()

# Get the API key from the environment
API_KEY = os.environ.get("GEMINI_API_KEY")

if not API_KEY:
    print("ERROR: GEMINI_API_KEY not found in .env file or environment.")
    sys.exit(1)

# Input / output paths (relative to the script's directory)
INPUT_FILE = "beszamolo.pdf"
OUTPUT_FILE = "beszamolo_styled_pptx.pptx"

# Gemini model to use
MODEL_NAME = "gemini-2.5-flash"

# Configuration Options
# Set to an integer (e.g. 5) to limit the number of generated slides, or None for unlimited.
MAX_SLIDES = 10

# Set to "light", "dark", or "none"
THEME = "light"


# ===========================================================================
# Step 1 — Configure the Gemini API
# ===========================================================================

def configure_api(api_key: str):
    """
    Configure the google-genai client with the provided API key.
    """
    client = genai.Client(api_key=api_key)
    print("[OK] Gemini API configured successfully.")
    return client


# ===========================================================================
# Step 2 — Read the input text file
# ===========================================================================

def read_input_text(filepath: str) -> str:
    """
    Extract text from a PDF file.

    Args:
        filepath: Path to the PDF file.

    Returns:
        The extracted text as a string.
    """
    script_dir = os.path.dirname(os.path.abspath(__file__))
    full_path = os.path.join(script_dir, filepath)

    try:
        reader = PdfReader(full_path)
        text_parts = []
        for page in reader.pages:
            text_parts.append(page.extract_text() or "")
        text = "\n".join(text_parts)
    except FileNotFoundError:
        print(f"ERROR: Input file not found: {full_path}")
        sys.exit(1)
    except Exception as e:
        print(f"ERROR: Could not read file '{full_path}': {e}")
        sys.exit(1)

    if not text.strip():
        print(f"ERROR: Extracted text is empty from: {full_path}")
        sys.exit(1)

    print(f"[OK] Read {len(text)} characters from '{filepath}'.")
    return text


# ===========================================================================
# Step 3 — Send the text to Gemini and request structured JSON slides
# ===========================================================================

def generate_slides(text: str, client, max_slides: int | None = None) -> str:
    """
    Send the input text to Gemini with a strict prompt that requests
    a JSON array of slide strings.

    Uses response_mime_type="application/json" to force the model to
    return valid JSON.

    Args:
        text: The source text to summarize into slides.
        client: The configured genai client.
        max_slides: Optional maximum number of slides to generate.

    Returns:
        The raw response text (expected to be a JSON array of strings).
    """
    # Build the prompt — instruct the model clearly
    prompt = (
        "You are a presentation designer. "
        "Read the following text carefully and extract the key points. "
        "Organize them into presentation slides. "
        "Each slide should have a short title line followed by 2-4 bullet points. "
        "Use the SAME language as the input text. "
    )

    if max_slides and isinstance(max_slides, int) and max_slides > 0:
        prompt += (
            f"CRITICAL INSTRUCTION: You MUST condense the entire text into EXACTLY {max_slides} slides. "
            "Do not cut off the story or information; ensure the presentation has a beginning, middle, and logical conclusion "
            f"within this strict {max_slides}-slide limit. "
        )

    prompt += (
        "Return your answer as a JSON array of strings, where each string is "
        "one slide's complete content (title + bullet points). "
        "Do NOT wrap the JSON in markdown code fences. "
        "Example format: "
        '[\"Slide Title\\n- Point 1\\n- Point 2\", \"Another Slide\\n- Point A\\n- Point B\"]\n\n'
        "Here is the text:\n\n"
        f"{text}"
    )

    # Configure generation to force JSON output
    config = types.GenerateContentConfig(
        response_mime_type="application/json",
        temperature=0.3,  # Lower temperature for more deterministic output
    )

    print(f"[...] Sending text to {MODEL_NAME} - this may take a few seconds...")

    try:
        response = client.models.generate_content(
            model=MODEL_NAME,
            contents=prompt,
            config=config,
        )
    except Exception as e:
        print(f"ERROR: Gemini API call failed: {e}")
        sys.exit(1)

    # Extract the text from the response
    raw_text = response.text
    print(f"[OK] Received response from Gemini ({len(raw_text)} characters).")
    return raw_text


# ===========================================================================
# Step 4 — Parse the JSON response into a list of slide strings
# ===========================================================================

def parse_slides(response_text: str) -> list[str]:
    """
    Parse the Gemini response as a JSON array of strings.

    Handles common issues:
    - Strips markdown code fences if the model wraps the JSON anyway.
    - Validates that the result is a list of strings.
    - Provides a clear error message on JSON decode failure.

    Args:
        response_text: Raw text from the Gemini response.

    Returns:
        A list of strings, each representing one slide's content.
    """
    # Strip potential markdown code fences (```json ... ```)
    cleaned = response_text.strip()
    if cleaned.startswith("```"):
        # Remove opening fence (e.g. ```json)
        cleaned = cleaned.split("\n", 1)[1] if "\n" in cleaned else cleaned[3:]
        # Remove closing fence
        if cleaned.endswith("```"):
            cleaned = cleaned[:-3].strip()

    # Attempt JSON parsing
    try:
        slides = json.loads(cleaned)
    except json.JSONDecodeError as e:
        print(f"ERROR: Failed to parse Gemini response as JSON: {e}")
        print("--- Raw response ---")
        print(response_text[:500])
        print("--- End of raw response ---")
        sys.exit(1)

    # Validate structure: must be a list of strings
    if not isinstance(slides, list):
        print(f"ERROR: Expected a JSON array, got {type(slides).__name__}.")
        sys.exit(1)

    # Coerce non-string items to strings (defensive)
    slides = [str(slide) for slide in slides]

    if len(slides) == 0:
        print("ERROR: Gemini returned an empty slide list.")
        sys.exit(1)

    print(f"[OK] Parsed {len(slides)} slides from the response.")
    return slides


# ===========================================================================
# Step 5 — Generate the PDF with one page per slide
# ===========================================================================

def create_pptx(slides: list[str], output_path: str) -> None:
    """
    Create a PPTX file with one slide per extracted slide content.

    Args:
        slides: List of slide content strings.
        output_path: File path for the generated PPTX.
    """
    script_dir = os.path.dirname(os.path.abspath(__file__))
    full_output = os.path.join(script_dir, output_path)

    prs = Presentation()

    # ---------------------------------------------------------------------------
    # Build slides
    # ---------------------------------------------------------------------------
    for i, slide_content in enumerate(slides, start=1):
        # Using a blank slide layout
        blank_slide_layout = prs.slide_layouts[6] 
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # --- Apply Background ---
        if THEME == "dark":
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(18, 18, 25)  # Very dark blue/grey
        elif THEME == "light" or THEME == "none":
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)  # White
        
        # Split content into title and body (assuming first line is title)
        lines = slide_content.strip().split("\n")
        title_text = lines[0] if lines else ""
        body_text = "\n".join(lines[1:]) if len(lines) > 1 else ""

        # --- Add Title Box ---
        title_box = slide.shapes.add_textbox(Pt(30), Pt(30), prs.slide_width - Pt(60), Pt(50))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Segoe UI"
        p.font.size = Pt(32)
        p.font.bold = True
        
        if THEME == "dark":
            p.font.color.rgb = RGBColor(0, 255, 204)  # Neon Cyan
        else:
            p.font.color.rgb = RGBColor(0, 0, 0)  # Black

        # --- Add Accent Line ---
        if THEME == "dark":
            accent_line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 
                Pt(32), Pt(85), 
                prs.slide_width - Pt(64), Pt(3)
            )
            accent_line.fill.solid()
            accent_line.fill.fore_color.rgb = RGBColor(138, 43, 226)  # Electric Purple
            accent_line.line.fill.background()  # Remove border
        elif THEME == "light":
            # Gradient blue line at the bottom using a single rectangle and OxmlElement
            accent_line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                0, prs.slide_height - Pt(30), 
                prs.slide_width, Pt(30)
            )
            accent_line.line.fill.background()  # Remove border
            
            # Explicitly set noFill so we know exactly where to insert gradFill in the XML sequence
            accent_line.fill.background()
            
            # Manipulate XML to add gradient in the correct sequence (before <a:ln>)
            spPr = accent_line._element.spPr
            noFill = spPr.find(qn('a:noFill'))
                
            gradFill = OxmlElement('a:gradFill')
            gsLst = OxmlElement('a:gsLst')
            
            # Stop 1: Dark Blue
            gs1 = OxmlElement('a:gs')
            gs1.set('pos', '0')
            srgbClr1 = OxmlElement('a:srgbClr')
            srgbClr1.set('val', '00008B')
            gs1.append(srgbClr1)
            gsLst.append(gs1)
            
            # Stop 2: Light Blue
            gs2 = OxmlElement('a:gs')
            gs2.set('pos', '100000')
            srgbClr2 = OxmlElement('a:srgbClr')
            srgbClr2.set('val', '87CEEB')
            gs2.append(srgbClr2)
            gsLst.append(gs2)
            
            gradFill.append(gsLst)
            
            # Linear gradient properties (0 degrees = left to right)
            lin = OxmlElement('a:lin')
            lin.set('ang', '0')
            lin.set('scaled', '1')
            gradFill.append(lin)
            
            # Insert gradFill right where noFill was, ensuring valid XML sequence
            if noFill is not None:
                noFill.addprevious(gradFill)
                spPr.remove(noFill)
            else:
                spPr.insert(1, gradFill)

        # --- Add Body Box ---
        if body_text:
            body_box = slide.shapes.add_textbox(Pt(30), Pt(100), prs.slide_width - Pt(60), prs.slide_height - Pt(130))
            tf_body = body_box.text_frame
            tf_body.word_wrap = True
            
            # Use the actual list layout capability in pptx by treating paragraphs properly
            is_first = True
            for line in body_text.split("\n"):
                if is_first:
                    p = tf_body.paragraphs[0]
                    is_first = False
                else:
                    p = tf_body.add_paragraph()
                
                # Check for bullet point format
                if line.strip().startswith("- ") or line.strip().startswith("* "):
                    p.text = line.strip()[2:]
                    p.level = 1  # This triggers the bullet point styling in pptx natively
                else:
                    p.text = line.strip()
                    p.level = 0
                
                # Style body text
                p.font.name = "Segoe UI"
                p.font.size = Pt(18)
                if THEME == "dark":
                    p.font.color.rgb = RGBColor(230, 230, 240)  # Off-white / light grey
                else:
                    p.font.color.rgb = RGBColor(0, 0, 0)  # Black
                    
                # Add a bit of space before paragraphs
                p.space_before = Pt(8)

    # ---------------------------------------------------------------------------
    # Write the PPTX to disk
    # ---------------------------------------------------------------------------
    try:
        prs.save(full_output)
    except OSError as e:
        print(f"ERROR: Could not write PPTX to '{full_output}': {e}")
        sys.exit(1)

    print(f"[OK] PPTX saved to '{full_output}' ({len(slides)} slides).")


# ===========================================================================
# Main pipeline
# ===========================================================================

def main() -> None:
    """
    Orchestrate the full pipeline:
        read text → send to Gemini → parse JSON slides → generate PDF
    """
    print("=" * 60)
    print("  PDF Presentation Generator")
    print("=" * 60)
    print()

    # 1. Configure the API
    client = configure_api(API_KEY)

    # 2. Read the source text
    text = read_input_text(INPUT_FILE)

    # 3. Generate slides via Gemini
    raw_response = generate_slides(text, client, MAX_SLIDES)

    # 4. Parse the JSON response
    slides = parse_slides(raw_response)

    # 5. Create the PPTX
    create_pptx(slides, OUTPUT_FILE)

    print()
    print("=" * 60)
    print(f"  Done! Open {OUTPUT_FILE} to view your presentation.")
    print("=" * 60)


if __name__ == "__main__":
    main()
